"""PPTX parser producing the same ordered, citable block contract as DOCX/PDF."""

from __future__ import annotations

import logging
import os
from collections.abc import Iterable
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..models import ContentBlock, ImageAsset
from .image_assets import image_asset_from_bytes
from .rasterizer import render_presentation_slides

logger = logging.getLogger(__name__)

PPTX_MEDIA_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
IMAGE_PLACEHOLDER = "[image]"


def parse_pptx(file_path: str, doc_id: str) -> list[ContentBlock]:
    """Parse slide text/tables and retain one faithful visual per slide.

    LibreOffice rendering preserves composite layout, charts, SmartArt, and
    drawing shapes. If rendering is unavailable, embedded picture bytes are
    retained instead; text and table extraction never depends on rendering.
    """
    _validate_file_path(file_path)
    presentation = Presentation(file_path)
    slide_renders = render_presentation_slides(file_path)
    blocks: list[ContentBlock] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title = _shape_text(title_shape) if title_shape is not None else ""
        heading_stack = [title] if title else []
        shape_records = list(
            _ordered_shapes(
                slide.shapes,
                slide_width=presentation.slide_width,
                slide_height=presentation.slide_height,
            )
        )

        if title:
            title_record = next(
                (record for record in shape_records if record[0] is title_shape),
                None,
            )
            blocks.append(
                _make_block(
                    doc_id=doc_id,
                    block_type="heading",
                    content=title,
                    heading_stack=heading_stack,
                    structural_meta={
                        "slide": slide_number,
                        "heading_level": 1,
                        **(title_record[2] if title_record else {}),
                    },
                    style_hint={"source": "pptx_title"},
                )
            )

        for shape, shape_path, position in shape_records:
            if shape is title_shape:
                continue
            common_meta = {
                "slide": slide_number,
                "shape_path": shape_path,
                **position,
            }
            if bool(getattr(shape, "has_table", False)):
                blocks.extend(
                    _table_blocks(
                        shape.table,
                        doc_id=doc_id,
                        heading_stack=heading_stack,
                        structural_meta=common_meta,
                    )
                )
                continue
            if bool(getattr(shape, "has_text_frame", False)):
                blocks.extend(
                    _text_blocks(
                        shape.text_frame,
                        doc_id=doc_id,
                        heading_stack=heading_stack,
                        structural_meta=common_meta,
                    )
                )

        notes = _speaker_notes(slide)
        if notes:
            blocks.append(
                _make_block(
                    doc_id=doc_id,
                    block_type="paragraph",
                    content=notes,
                    heading_stack=heading_stack,
                    structural_meta={"slide": slide_number, "speaker_notes": True},
                    style_hint={"source": "pptx_speaker_notes"},
                )
            )

        rendered = slide_renders.get(slide_number)
        if rendered:
            asset = image_asset_from_bytes(rendered, PPTX_MEDIA_TYPE)
            if asset is not None:
                blocks.append(
                    _image_block(
                        doc_id=doc_id,
                        image=asset,
                        heading_stack=heading_stack,
                        structural_meta={
                            "slide": slide_number,
                            "visual_scope": "full_slide",
                        },
                        source="pptx_slide_render",
                    )
                )
        else:
            blocks.extend(
                _embedded_picture_blocks(
                    shape_records,
                    doc_id=doc_id,
                    slide_number=slide_number,
                    heading_stack=heading_stack,
                )
            )

    image_index = 0
    for ordinal, block in enumerate(blocks):
        block.ordinal = ordinal
        block.id = f"{doc_id}/b-{ordinal:04d}"
        if block.block_type == "image":
            block.structural_meta["image_index"] = image_index
            image_index += 1
    return blocks


def _validate_file_path(file_path: str) -> None:
    if not isinstance(file_path, str) or not file_path:
        raise ValueError("file_path must be a non-empty string")
    if not file_path.lower().endswith(".pptx"):
        raise ValueError("file_path must point to a .pptx file")
    if not os.path.exists(file_path):
        raise FileNotFoundError(file_path)


def _ordered_shapes(
    shapes: Iterable[Any],
    *,
    slide_width: int,
    slide_height: int,
    prefix: str = "",
) -> Iterable[tuple[Any, str, dict[str, float]]]:
    indexed = list(enumerate(shapes))
    indexed.sort(
        key=lambda item: (
            int(getattr(item[1], "top", 0) or 0),
            int(getattr(item[1], "left", 0) or 0),
            item[0],
        )
    )
    for original_index, shape in indexed:
        path = f"{prefix}.{original_index}" if prefix else str(original_index)
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from _ordered_shapes(
                shape.shapes,
                slide_width=slide_width,
                slide_height=slide_height,
                prefix=path,
            )
            continue
        yield shape, path, _normalized_position(shape, slide_width, slide_height)


def _normalized_position(
    shape: Any, slide_width: int, slide_height: int
) -> dict[str, float]:
    width = max(int(slide_width), 1)
    height = max(int(slide_height), 1)
    return {
        "x": round(int(getattr(shape, "left", 0) or 0) / width, 4),
        "y": round(int(getattr(shape, "top", 0) or 0) / height, 4),
        "width": round(int(getattr(shape, "width", 0) or 0) / width, 4),
        "height": round(int(getattr(shape, "height", 0) or 0) / height, 4),
    }


def _shape_text(shape: Any) -> str:
    if shape is None or not bool(getattr(shape, "has_text_frame", False)):
        return ""
    return "\n".join(
        paragraph.text.strip()
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text.strip()
    )


def _text_blocks(
    text_frame: Any,
    *,
    doc_id: str,
    heading_stack: list[str],
    structural_meta: dict[str, Any],
) -> list[ContentBlock]:
    blocks: list[ContentBlock] = []
    for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
        text = paragraph.text.strip()
        if not text:
            continue
        blocks.append(
            _make_block(
                doc_id=doc_id,
                block_type="paragraph",
                content=text,
                heading_stack=heading_stack,
                structural_meta={
                    **structural_meta,
                    "paragraph_index": paragraph_index,
                    "bullet_level": int(getattr(paragraph, "level", 0) or 0),
                },
                style_hint={"source": "pptx_text"},
            )
        )
    return blocks


def _table_blocks(
    table: Any,
    *,
    doc_id: str,
    heading_stack: list[str],
    structural_meta: dict[str, Any],
) -> list[ContentBlock]:
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    rows = [row for row in rows if any(row)]
    if not rows:
        return []
    width = max(len(row) for row in rows)
    normalized = [_normalize_row(row, width) for row in rows]

    if width == 1:
        return [
            _make_block(
                doc_id=doc_id,
                block_type="paragraph",
                content=row[0],
                heading_stack=heading_stack,
                structural_meta={**structural_meta, "row_index": row_index},
                style_hint={"source": "pptx_single_column_table"},
            )
            for row_index, row in enumerate(normalized)
            if row[0]
        ]

    headers = normalized[0]
    data_rows = normalized[1:]
    if not data_rows:
        content = " | ".join(value for value in headers if value)
        return [
            _make_block(
                doc_id=doc_id,
                block_type="paragraph",
                content=content,
                heading_stack=heading_stack,
                structural_meta={
                    **structural_meta,
                    "row_index": 0,
                    "column_headers": headers,
                },
                style_hint={"source": "pptx_table_headers"},
            )
        ] if content else []

    blocks: list[ContentBlock] = []
    for row_index, row in enumerate(data_rows, start=1):
        content = ", ".join(
            f"{headers[index]}: {value}" if headers[index] else value
            for index, value in enumerate(row)
            if value
        )
        if not content:
            continue
        blocks.append(
            _make_block(
                doc_id=doc_id,
                block_type="table_row",
                content=content,
                heading_stack=heading_stack,
                structural_meta={
                    **structural_meta,
                    "row_index": row_index,
                    "column_headers": headers,
                },
                style_hint={"source": "pptx_table_row"},
            )
        )
    return blocks


def _normalize_row(row: list[str], width: int) -> list[str]:
    return row[:width] + [""] * max(0, width - len(row))


def _speaker_notes(slide: Any) -> str:
    if not bool(getattr(slide, "has_notes_slide", False)):
        return ""
    try:
        text_frame = slide.notes_slide.notes_text_frame
    except (AttributeError, ValueError):
        return ""
    return "\n".join(
        paragraph.text.strip()
        for paragraph in text_frame.paragraphs
        if paragraph.text.strip()
    )


def _embedded_picture_blocks(
    shape_records: list[tuple[Any, str, dict[str, float]]],
    *,
    doc_id: str,
    slide_number: int,
    heading_stack: list[str],
) -> list[ContentBlock]:
    blocks: list[ContentBlock] = []
    for shape, shape_path, position in shape_records:
        try:
            picture = shape.image
            data = picture.blob
            source_media_type = picture.content_type
        except (AttributeError, ValueError, KeyError):
            continue
        asset = image_asset_from_bytes(data, source_media_type)
        if asset is None:
            logger.warning(
                "Could not retain PPTX picture on slide %s at shape %s",
                slide_number,
                shape_path,
            )
            continue
        blocks.append(
            _image_block(
                doc_id=doc_id,
                image=asset,
                heading_stack=heading_stack,
                structural_meta={
                    "slide": slide_number,
                    "shape_path": shape_path,
                    "visual_scope": "embedded_picture",
                    **position,
                },
                source="pptx_picture",
            )
        )
    return blocks


def _image_block(
    *,
    doc_id: str,
    image: ImageAsset,
    heading_stack: list[str],
    structural_meta: dict[str, Any],
    source: str,
) -> ContentBlock:
    return ContentBlock(
        id="",
        doc_id=doc_id,
        ordinal=-1,
        block_type="image",
        content=IMAGE_PLACEHOLDER,
        heading_stack=heading_stack.copy(),
        structural_meta=structural_meta,
        style_hint={"source": source},
        image=image,
    )


def _make_block(
    *,
    doc_id: str,
    block_type: str,
    content: str,
    heading_stack: list[str],
    structural_meta: dict[str, Any],
    style_hint: dict[str, Any],
) -> ContentBlock:
    return ContentBlock(
        id="",
        doc_id=doc_id,
        ordinal=-1,
        block_type=block_type,
        content=content,
        heading_stack=heading_stack.copy(),
        structural_meta=structural_meta,
        style_hint=style_hint,
    )
