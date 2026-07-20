from __future__ import annotations

import base64
import tempfile
import unittest
from io import BytesIO
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches

from services.chunker import run_pipeline
from services.chunker.stages.rasterizer import render_presentation_slides


PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


def build_presentation(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Malaria vaccine profile"

    text_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(5), Inches(1.2)
    )
    text_box.text_frame.text = "Target population: children"
    second = text_box.text_frame.add_paragraph()
    second.text = "Two-dose primary series"
    second.level = 1

    table = slide.shapes.add_table(
        2, 2, Inches(0.8), Inches(3), Inches(5), Inches(1.2)
    ).table
    table.cell(0, 0).text = "Measure"
    table.cell(0, 1).text = "Target"
    table.cell(1, 0).text = "Efficacy"
    table.cell(1, 1).text = ">= 75%"

    slide.shapes.add_picture(
        BytesIO(PNG_1X1), Inches(6.2), Inches(1.5), width=Inches(1)
    )
    presentation.save(path)


class ChunkerPptxTests(unittest.TestCase):
    def test_slide_renderer_emits_png_when_office_runtime_is_available(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "profile.pptx"
            build_presentation(path)
            rendered = render_presentation_slides(str(path))

        if not rendered:
            self.skipTest("LibreOffice presentation rendering is unavailable")
        self.assertEqual(set(rendered), {1})
        self.assertTrue(rendered[1].startswith(b"\x89PNG\r\n\x1a\n"))

    def test_pptx_emits_text_table_and_one_rendered_slide_visual(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "profile.pptx"
            build_presentation(path)
            with patch(
                "services.chunker.stages.parser_pptx.render_presentation_slides",
                return_value={1: PNG_1X1},
            ):
                blocks = run_pipeline(str(path), "profile")

        self.assertEqual(
            [block.id for block in blocks],
            [f"profile/b-{index:04d}" for index in range(len(blocks))],
        )
        self.assertEqual(blocks[0].block_type, "heading")
        self.assertEqual(blocks[0].content, "Malaria vaccine profile")
        self.assertTrue(
            any(block.content == "Target population: children" for block in blocks)
        )
        self.assertTrue(
            any(
                block.block_type == "table_row"
                and block.content == "Measure: Efficacy, Target: >= 75%"
                for block in blocks
            )
        )
        image_blocks = [block for block in blocks if block.block_type == "image"]
        self.assertEqual(len(image_blocks), 1)
        image = image_blocks[0]
        self.assertEqual(image.structural_meta["slide"], 1)
        self.assertEqual(image.structural_meta["visual_scope"], "full_slide")
        self.assertIsNotNone(image.image)
        assert image.image is not None
        self.assertEqual(image.image.media_type, "image/png")
        self.assertIn("presentationml.presentation", image.image.source_media_type)

    def test_pptx_retains_embedded_picture_when_slide_rendering_is_unavailable(self) -> None:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "profile.pptx"
            build_presentation(path)
            with patch(
                "services.chunker.stages.parser_pptx.render_presentation_slides",
                return_value={},
            ):
                blocks = run_pipeline(str(path), "profile")

        image_blocks = [block for block in blocks if block.block_type == "image"]
        self.assertEqual(len(image_blocks), 1)
        image = image_blocks[0]
        self.assertEqual(image.structural_meta["visual_scope"], "embedded_picture")
        self.assertIsNotNone(image.image)
        assert image.image is not None
        self.assertEqual(image.image.source_media_type, "image/png")


if __name__ == "__main__":
    unittest.main()
